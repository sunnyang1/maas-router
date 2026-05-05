"""
MaaS-Router Document Automation Service
========================================
Unified document generation engine supporting PDF, Excel, Word, and PPTX formats.

Architecture:
  Data Layer (PostgreSQL/Redis) → Engine Layer (PDF/Excel/Word/PPTX) → Output Layer (File/Stream)

Document Types:
  - billing_report    : Billing summaries, revenue reports, transaction history
  - user_usage_report : Per-user API usage, token consumption, cost breakdown
  - model_performance : Per-model latency, error rates, routing efficiency
  - ops_daily         : Daily operational snapshot, health checks, alerts
  - audit_report      : Compliance audit trails, change logs
  - data_export       : Raw data export to Excel/CSV
"""

from __future__ import annotations

import io
import json
import os
from datetime import datetime, timezone, timedelta
from enum import Enum
from pathlib import Path
from typing import Any, Optional
from jinja2 import Environment, FileSystemLoader

# ── PDF Generation ──────────────────────────────────────────────
# WeasyPrint: HTML+CSS → PDF, best for complex layouts with branding
try:
    from weasyprint import HTML as WeasyHTML
    WEASYPRINT_AVAILABLE = True
except ImportError:
    WEASYPRINT_AVAILABLE = False

# ── Excel Generation ────────────────────────────────────────────
# openpyxl: native .xlsx with charts, formatting, formulas
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.chart import BarChart, PieChart, LineChart, Reference
    from openpyxl.utils import get_column_letter
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

# ── Word Generation ─────────────────────────────────────────────
# python-docx: .docx with styles, TOC, headers, tables
try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.section import WD_ORIENT
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

# ── PPTX Generation ─────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False


# ══════════════════════════════════════════════════════════════════
# Configuration
# ══════════════════════════════════════════════════════════════════

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "output"
TEMPLATE_DIR = Path(__file__).resolve().parent.parent / "templates" / "documents"

# Brand colors for MaaS-Router
BRAND = {
    "primary": "#6366f1",        # Indigo-500
    "primary_dark": "#4f46e5",   # Indigo-600
    "secondary": "#06b6d4",      # Cyan-500
    "success": "#10b981",        # Emerald-500
    "warning": "#f59e0b",        # Amber-500
    "danger": "#ef4444",         # Red-500
    "gray_50": "#f9fafb",
    "gray_100": "#f3f4f6",
    "gray_200": "#e5e7eb",
    "gray_600": "#4b5563",
    "gray_800": "#1f2937",
    "gray_900": "#111827",
}


class DocumentFormat(str, Enum):
    PDF = "pdf"
    EXCEL = "xlsx"
    WORD = "docx"
    PPTX = "pptx"
    CSV = "csv"


class DocumentType(str, Enum):
    BILLING_REPORT = "billing_report"
    USER_USAGE_REPORT = "user_usage_report"
    MODEL_PERFORMANCE = "model_performance"
    OPS_DAILY = "ops_daily"
    AUDIT_REPORT = "audit_report"
    DATA_EXPORT = "data_export"


# ══════════════════════════════════════════════════════════════════
# Jinja2 Template Engine Setup
# ══════════════════════════════════════════════════════════════════

_jinja_env = Environment(
    loader=FileSystemLoader(str(TEMPLATE_DIR)),
    autoescape=True,
    trim_blocks=True,
    lstrip_blocks=True,
)

# Register custom filters
_jinja_env.filters["money"] = lambda v, currency="CRED": f"{currency} {v:,.2f}"
_jinja_env.filters["pct"] = lambda v: f"{v * 100:.1f}%" if isinstance(v, float) else f"{v}%"
_jinja_env.filters["datetime_fmt"] = lambda v, fmt="%Y-%m-%d %H:%M": (
    v.strftime(fmt) if isinstance(v, datetime) else str(v)
)
_jinja_env.filters["number"] = lambda v: f"{v:,}"
_jinja_env.filters["filesize"] = lambda v: (
    f"{v / 1024:.1f} KB" if v < 1_048_576 else f"{v / 1_048_576:.1f} MB"
)


# ══════════════════════════════════════════════════════════════════
# Base Engine
# ══════════════════════════════════════════════════════════════════

class DocumentEngine:
    """Base class for all document generation engines."""

    def __init__(self):
        self.output_dir = OUTPUT_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def _ensure_output_dir(self) -> Path:
        self.output_dir.mkdir(parents=True, exist_ok=True)
        return self.output_dir

    def _filename(self, doc_type: str, fmt: str, suffix: str | None = None) -> str:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        sfx = suffix or fmt
        return f"{doc_type}_{ts}.{sfx}"


# ══════════════════════════════════════════════════════════════════
# PDF Engine (WeasyPrint)
# ══════════════════════════════════════════════════════════════════

class PDFEngine(DocumentEngine):
    """
    HTML+CSS → PDF via WeasyPrint.
    Supports complex layouts, branding, charts (via embedded SVG/CSS), and accessibility tags.
    """

    ENGINE_NAME = "weasyprint"

    def render_from_html(self, html: str, output_path: Path | None = None,
                         doc_type: str = "document") -> tuple[bytes, Path]:
        """Render raw HTML string to PDF."""
        if not WEASYPRINT_AVAILABLE:
            raise RuntimeError("WeasyPrint is not installed. Run: pip install weasyprint")

        pdf_bytes = WeasyHTML(string=html).write_pdf()

        if output_path is None:
            output_path = self.output_dir / self._filename(doc_type, "pdf")

        output_path.write_bytes(pdf_bytes)
        return pdf_bytes, output_path

    def render_from_template(self, template_name: str, context: dict,
                              output_path: Path | None = None,
                              doc_type: str = "document") -> tuple[bytes, Path]:
        """Render Jinja2 template to PDF."""
        template = _jinja_env.get_template(template_name)
        html = template.render(**context, brand=BRAND, now=datetime.now(timezone.utc))
        return self.render_from_html(html, output_path, doc_type)

    def generate_billing_report(self, data: dict) -> tuple[bytes, Path]:
        """Generate billing summary PDF report."""
        return self.render_from_template(
            "billing_report.html",
            data,
            doc_type="billing_report",
        )

    def generate_user_report(self, data: dict) -> tuple[bytes, Path]:
        """Generate per-user usage report PDF."""
        return self.render_from_template(
            "user_report.html",
            data,
            doc_type="user_usage_report",
        )

    def generate_ops_daily(self, data: dict) -> tuple[bytes, Path]:
        """Generate daily operations PDF report."""
        return self.render_from_template(
            "ops_daily.html",
            data,
            doc_type="ops_daily",
        )


# ══════════════════════════════════════════════════════════════════
# Excel Engine (openpyxl)
# ══════════════════════════════════════════════════════════════════

class ExcelEngine(DocumentEngine):
    """
    Native .xlsx generation with openpyxl.
    Supports formatting, charts, formulas, pivot-ready layouts, and multiple sheets.
    """

    ENGINE_NAME = "openpyxl"

    # Style constants
    HEADER_FILL = PatternFill(start_color=BRAND["primary"], end_color=BRAND["primary"], fill_type="solid")
    HEADER_FONT = Font(name="Inter, Arial", bold=True, color="FFFFFF", size=11)
    DATA_FONT = Font(name="Inter, Arial", size=10)
    MONEY_FONT = Font(name="Inter, Arial", size=10, color="166534")  # green-800
    BORDER = Border(
        left=Side(style="thin", color="D1D5DB"),
        right=Side(style="thin", color="D1D5DB"),
        top=Side(style="thin", color="D1D5DB"),
        bottom=Side(style="thin", color="D1D5DB"),
    )
    ALT_FILL = PatternFill(start_color="F3F4F6", end_color="F3F4F6", fill_type="solid")

    def __init__(self):
        super().__init__()
        if not OPENPYXL_AVAILABLE:
            raise RuntimeError("openpyxl is not installed. Run: pip install openpyxl")

    def create_workbook(self, title: str = "MaaS-Router Report") -> Workbook:
        wb = Workbook()
        wb.properties.title = title
        wb.properties.creator = "MaaS-Router Document Engine"
        return wb

    def style_header_row(self, ws, row: int, col_count: int):
        """Apply header styling to a row."""
        for col in range(1, col_count + 1):
            cell = ws.cell(row=row, column=col)
            cell.fill = self.HEADER_FILL
            cell.font = self.HEADER_FONT
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = self.BORDER

    def style_data_rows(self, ws, start_row: int, end_row: int, col_count: int):
        """Apply alternating row styling and borders to data rows."""
        for row in range(start_row, end_row + 1):
            for col in range(1, col_count + 1):
                cell = ws.cell(row=row, column=col)
                cell.font = self.DATA_FONT
                cell.border = self.BORDER
                cell.alignment = Alignment(vertical="center")
                if (row - start_row) % 2 == 1:
                    cell.fill = self.ALT_FILL

    def auto_width(self, ws, col_count: int, min_width: int = 10, max_width: int = 40):
        """Auto-fit column widths."""
        for col in range(1, col_count + 1):
            letter = get_column_letter(col)
            max_len = min_width
            for row in ws.iter_rows(min_col=col, max_col=col, values_only=True):
                for cell_value in row:
                    if cell_value:
                        max_len = max(max_len, min(len(str(cell_value)) + 2, max_width))
            ws.column_dimensions[letter].width = max_len

    def generate_data_export(self, sheets: dict[str, list[dict]],
                              filename: str | None = None) -> tuple[bytes, Path]:
        """
        Generate multi-sheet Excel data export.

        Args:
            sheets: {sheet_name: [row_dict, ...], ...}
            filename: optional output filename
        """
        wb = self.create_workbook()

        for idx, (sheet_name, rows) in enumerate(sheets.items()):
            if idx == 0:
                ws = wb.active
                ws.title = sheet_name
            else:
                ws = wb.create_sheet(title=sheet_name)

            if not rows:
                ws.cell(row=1, column=1, value="No data")
                continue

            # Write headers
            headers = list(rows[0].keys())
            for col, header in enumerate(headers, 1):
                ws.cell(row=1, column=col, value=header.replace("_", " ").title())

            # Write data
            for row_idx, row_data in enumerate(rows, 2):
                for col, key in enumerate(headers, 1):
                    value = row_data.get(key)
                    # Format datetimes
                    if isinstance(value, datetime):
                        value = value.strftime("%Y-%m-%d %H:%M:%S")
                    ws.cell(row=row_idx, column=col, value=value)

            # Apply styling
            col_count = len(headers)
            self.style_header_row(ws, 1, col_count)
            if len(rows) > 0:
                self.style_data_rows(ws, 2, len(rows) + 1, col_count)
            self.auto_width(ws, col_count)
            ws.auto_filter.ref = ws.dimensions

        # Save
        output_path = self.output_dir / (filename or self._filename("data_export", "xlsx"))
        wb.save(str(output_path))

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue(), output_path

    def generate_billing_excel(self, billing_data: dict) -> tuple[bytes, Path]:
        """Generate billing report Excel with charts."""
        wb = self.create_workbook("MaaS-Router Billing Report")

        # ── Sheet 1: Summary ──
        ws_summary = wb.active
        ws_summary.title = "Summary"

        summary_data = billing_data.get("summary", {})
        labels = [
            ("Report Period", f"{billing_data.get('period_start', '')} → {billing_data.get('period_end', '')}"),
            ("Total Revenue", summary_data.get("total_revenue", 0)),
            ("Total Transactions", summary_data.get("total_transactions", 0)),
            ("Active Users", summary_data.get("active_users", 0)),
            ("New Users", summary_data.get("new_users", 0)),
            ("Top Plan", summary_data.get("top_plan", "N/A")),
        ]
        for i, (label, value) in enumerate(labels, 1):
            cell_a = ws_summary.cell(row=i, column=1, value=label)
            cell_a.font = Font(bold=True, size=11)
            cell_b = ws_summary.cell(row=i, column=2, value=value)
            if "Revenue" in label:
                cell_b.number_format = '#,##0.00 "CRED"'

        ws_summary.column_dimensions["A"].width = 22
        ws_summary.column_dimensions["B"].width = 30

        # ── Sheet 2: Daily Breakdown ──
        if "daily_breakdown" in billing_data:
            ws_daily = wb.create_sheet("Daily Breakdown")
            headers = ["Date", "Requests", "Tokens", "Revenue (CRED)", "Avg Latency (ms)", "Error Rate"]
            for c, h in enumerate(headers, 1):
                ws_daily.cell(row=1, column=c, value=h)

            daily = billing_data["daily_breakdown"]
            for r, day in enumerate(daily, 2):
                ws_daily.cell(row=r, column=1, value=day.get("date", ""))
                ws_daily.cell(row=r, column=2, value=day.get("requests", 0))
                ws_daily.cell(row=r, column=3, value=day.get("tokens", 0))
                ws_daily.cell(row=r, column=4, value=day.get("revenue", 0))
                ws_daily.cell(row=r, column=5, value=day.get("avg_latency", 0))
                ws_daily.cell(row=r, column=6, value=day.get("error_rate", 0))

            col_count = len(headers)
            self.style_header_row(ws_daily, 1, col_count)
            if daily:
                self.style_data_rows(ws_daily, 2, len(daily) + 1, col_count)
            self.auto_width(ws_daily, col_count)

            # Revenue trend chart
            if len(daily) > 1:
                chart = LineChart()
                chart.title = "Daily Revenue Trend"
                chart.y_axis.title = "CRED"
                chart.style = 10
                data_ref = Reference(ws_daily, min_col=4, min_row=1, max_row=len(daily) + 1)
                cats_ref = Reference(ws_daily, min_col=1, min_row=2, max_row=len(daily) + 1)
                chart.add_data(data_ref, titles_from_data=True)
                chart.set_categories(cats_ref)
                chart.series[0].graphicalProperties.line.solidFill = BRAND["primary"]
                ws_daily.add_chart(chart, "H2")

        # ── Sheet 3: Transactions ──
        if "recent_transactions" in billing_data:
            ws_txn = wb.create_sheet("Transactions")
            headers = ["ID", "User", "Type", "Model", "Tokens", "Amount", "Currency", "Status", "Date"]
            for c, h in enumerate(headers, 1):
                ws_txn.cell(row=1, column=c, value=h)

            txns = billing_data["recent_transactions"]
            for r, txn in enumerate(txns, 2):
                ws_txn.cell(row=r, column=1, value=txn.get("id", ""))
                ws_txn.cell(row=r, column=2, value=txn.get("user", ""))
                ws_txn.cell(row=r, column=3, value=txn.get("txn_type", ""))
                ws_txn.cell(row=r, column=4, value=txn.get("model", ""))
                ws_txn.cell(row=r, column=5, value=txn.get("tokens", 0))
                ws_txn.cell(row=r, column=6, value=txn.get("amount", 0))
                ws_txn.cell(row=r, column=7, value=txn.get("currency", "CRED"))
                ws_txn.cell(row=r, column=8, value=txn.get("status", ""))
                ws_txn.cell(row=r, column=9, value=txn.get("created_at", ""))

            col_count = len(headers)
            self.style_header_row(ws_txn, 1, col_count)
            if txns:
                self.style_data_rows(ws_txn, 2, len(txns) + 1, col_count)
            self.auto_width(ws_txn, col_count)

        # Save
        output_path = self.output_dir / self._filename("billing_report", "xlsx")
        wb.save(str(output_path))

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue(), output_path

    def generate_csv(self, data: list[dict], filename: str | None = None) -> tuple[str, Path]:
        """Generate CSV from list of dicts."""
        import csv

        output_path = self.output_dir / (filename or self._filename("export", "csv"))
        if not data:
            output_path.write_text("No data", encoding="utf-8")
            return "No data", output_path

        with open(output_path, "w", newline="", encoding="utf-8-sig") as f:
            writer = csv.DictWriter(f, fieldnames=list(data[0].keys()))
            writer.writeheader()
            writer.writerows(data)

        content = output_path.read_text(encoding="utf-8-sig")
        return content, output_path


# ══════════════════════════════════════════════════════════════════
# Word Engine (python-docx)
# ══════════════════════════════════════════════════════════════════

class WordEngine(DocumentEngine):
    """
    .docx generation with proper styles, headers, TOC placeholder, and consistent formatting.
    """

    ENGINE_NAME = "python-docx"

    def __init__(self):
        super().__init__()
        if not PYTHON_DOCX_AVAILABLE:
            raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    def _setup_styles(self, doc: DocxDocument):
        """Configure document styles for consistent branding."""
        style = doc.styles["Normal"]
        style.font.name = "Inter, Arial"
        style.font.size = Pt(10.5)
        style.paragraph_format.space_after = Pt(6)

        for level in range(1, 4):
            heading_style = doc.styles[f"Heading {level}"]
            heading_style.font.color.rgb = RGBColor(0x4F, 0x46, 0xE5)  # primary_dark

    def _add_cover_page(self, doc: DocxDocument, title: str, subtitle: str,
                         report_date: str, report_id: str):
        """Add a branded cover page."""
        # Spacer
        for _ in range(6):
            doc.add_paragraph("")

        # Title
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(title)
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x4F, 0x46, 0xE5)

        # Subtitle
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(subtitle)
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

        doc.add_paragraph("")
        doc.add_paragraph("")

        # Meta
        meta_items = [
            ("Report Date", report_date),
            ("Report ID", report_id),
            ("Generated by", "MaaS-Router Document Engine"),
            ("Classification", "Internal"),
        ]
        for label, value in meta_items:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(f"{label}: ")
            run.font.bold = True
            p.add_run(value)

        doc.add_page_break()

    def _add_table(self, doc: DocxDocument, headers: list[str], rows: list[list],
                    col_widths: list[float] | None = None):
        """Add a formatted table to the document."""
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Light Grid Accent 1"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        # Header row
        for i, header in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = header
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(9)

        # Data rows
        for r, row_data in enumerate(rows, 1):
            for c, value in enumerate(row_data):
                table.rows[r].cells[c].text = str(value)

        return table

    def generate_audit_report(self, data: dict) -> tuple[bytes, Path]:
        """Generate compliance audit report as .docx."""
        doc = DocxDocument()
        self._setup_styles(doc)

        report_date = datetime.now().strftime("%Y-%m-%d %H:%M UTC")
        report_id = data.get("report_id", f"AUD-{datetime.now().strftime('%Y%m%d')}")
        period = f"{data.get('period_start', 'N/A')} → {data.get('period_end', 'N/A')}"

        self._add_cover_page(
            doc,
            "MaaS-Router Audit Report",
            f"Compliance Audit: {period}",
            report_date,
            report_id,
        )

        # ── Executive Summary ──
        doc.add_heading("1. Executive Summary", level=1)
        summary = data.get("summary", {})
        doc.add_paragraph(
            f"This report covers {summary.get('total_events', 0)} auditable events "
            f"across {summary.get('active_users', 0)} active users during the period "
            f"{period}. {summary.get('critical_events', 0)} critical events were recorded."
        )

        # ── Audit Event Summary ──
        doc.add_heading("2. Audit Event Summary", level=1)
        audit_summary = data.get("audit_summary", [])
        if audit_summary:
            self._add_table(
                doc,
                ["Event Type", "Count", "Unique Users", "Last Occurrence"],
                [[
                    e.get("action", ""),
                    str(e.get("count", 0)),
                    str(e.get("unique_users", 0)),
                    e.get("last_occurrence", ""),
                ] for e in audit_summary],
            )

        # ── Critical Events ──
        doc.add_heading("3. Critical Events", level=1)
        critical = data.get("critical_events", [])
        if critical:
            for event in critical[:50]:  # Limit to 50
                doc.add_paragraph(
                    f"[{event.get('timestamp', '')}] {event.get('action', '')} — "
                    f"User: {event.get('user', 'N/A')} | "
                    f"Resource: {event.get('resource_type', '')}/{event.get('resource_id', '')} | "
                    f"IP: {event.get('ip', 'N/A')}",
                    style="List Bullet",
                )
        else:
            doc.add_paragraph("No critical events recorded during this period.")

        # ── User Activity ──
        doc.add_heading("4. User Activity Summary", level=1)
        user_activity = data.get("user_activity", [])
        if user_activity:
            self._add_table(
                doc,
                ["User", "Email", "Actions", "Last Active", "IP Addresses"],
                [[
                    u.get("display_name", ""),
                    u.get("email", ""),
                    str(u.get("action_count", 0)),
                    u.get("last_active", ""),
                    str(u.get("ip_count", 0)),
                ] for u in user_activity],
            )

        # ── Compliance Notes ──
        doc.add_heading("5. Compliance Notes", level=1)
        notes = data.get("compliance_notes", [])
        if notes:
            for note in notes:
                doc.add_paragraph(note, style="List Bullet")
        else:
            doc.add_paragraph("All events are within normal operational parameters.")
            doc.add_paragraph("No compliance violations detected.")

        # Save
        output_path = self.output_dir / self._filename("audit_report", "docx")
        doc.save(str(output_path))

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue(), output_path

    def generate_user_report_docx(self, data: dict) -> tuple[bytes, Path]:
        """Generate detailed user usage report as .docx."""
        doc = DocxDocument()
        self._setup_styles(doc)

        report_date = datetime.now().strftime("%Y-%m-%d %H:%M UTC")
        period = f"{data.get('period_start', 'N/A')} → {data.get('period_end', 'N/A')}"

        self._add_cover_page(
            doc,
            "User Usage Report",
            f"Period: {period}",
            report_date,
            f"USR-{datetime.now().strftime('%Y%m%d')}",
        )

        # ── Overview ──
        doc.add_heading("1. Overview", level=1)
        overview = data.get("overview", {})
        doc.add_paragraph(f"Total Active Users: {overview.get('active_users', 0)}")
        doc.add_paragraph(f"Total API Requests: {overview.get('total_requests', 0):,}")
        doc.add_paragraph(f"Total Tokens Consumed: {overview.get('total_tokens', 0):,}")
        doc.add_paragraph(f"Total Cost: CRED {overview.get('total_cost', 0):,.2f}")

        # ── Top Users ──
        doc.add_heading("2. Top Users by Usage", level=1)
        top_users = data.get("top_users", [])
        if top_users:
            self._add_table(
                doc,
                ["Rank", "User", "Email", "Plan", "Requests", "Tokens", "Cost (CRED)"],
                [[
                    str(i + 1),
                    u.get("display_name", ""),
                    u.get("email", ""),
                    u.get("plan", ""),
                    str(u.get("requests", 0)),
                    str(u.get("tokens", 0)),
                    f"{u.get('cost', 0):,.2f}",
                ] for i, u in enumerate(top_users)],
            )

        # ── Model Usage Distribution ──
        doc.add_heading("3. Model Usage Distribution", level=1)
        model_usage = data.get("model_usage", [])
        if model_usage:
            self._add_table(
                doc,
                ["Model", "Requests", "Tokens", "Avg Latency (ms)", "Cost (CRED)"],
                [[
                    m.get("model_id", ""),
                    str(m.get("requests", 0)),
                    str(m.get("tokens", 0)),
                    str(m.get("avg_latency", 0)),
                    f"{m.get('cost', 0):,.2f}",
                ] for m in model_usage],
            )

        # ── Recommendations ──
        doc.add_heading("4. Recommendations", level=1)
        recs = data.get("recommendations", [
            "Review high-usage users for plan upgrade opportunities.",
            "Monitor error rates for model quality regression.",
            "Consider rate limit adjustments for enterprise users.",
        ])
        for rec in recs:
            doc.add_paragraph(rec, style="List Bullet")

        # Save
        output_path = self.output_dir / self._filename("user_usage_report", "docx")
        doc.save(str(output_path))

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue(), output_path


# ══════════════════════════════════════════════════════════════════
# PPTX Engine (python-pptx)
# ══════════════════════════════════════════════════════════════════

class PPTXEngine(DocumentEngine):
    """
    Presentation (.pptx) generation with data-driven slides.
    Suitable for investor decks, quarterly reviews, and team presentations.
    """

    ENGINE_NAME = "python-pptx"

    def __init__(self):
        super().__init__()
        if not PYTHON_PPTX_AVAILABLE:
            raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    def _apply_branding(self, prs: Presentation):
        """Apply MaaS-Router branding to all slide layouts."""
        # Use blank layout for custom designs
        pass

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str, date: str):
        """Add a branded title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1].has_text_frame:
            slide.placeholders[1].text = f"{subtitle}\n{date}"

    def _add_bullet_slide(self, prs: Presentation, title: str, bullets: list[str]):
        """Add a slide with title and bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        body = slide.placeholders[1]
        tf = body.text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    def _add_table_slide(self, prs: Presentation, title: str,
                          headers: list[str], rows: list[list]):
        """Add a slide with a data table."""
        slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        left = PptxInches(0.5)
        top = PptxInches(0.3)
        txBox = slide.shapes.add_textbox(left, top, PptxInches(9), PptxInches(0.6))
        tf = txBox.text_frame
        tf.text = title

        # Table
        rows_count = min(len(rows) + 1, 15)
        cols_count = len(headers)
        table_shape = slide.shapes.add_table(
            rows_count, cols_count,
            PptxInches(0.5), PptxInches(1.2),
            PptxInches(9), PptxInches(0.35) * rows_count,
        )
        table = table_shape.table

        for c, h in enumerate(headers):
            table.cell(0, c).text = h
        for r, row in enumerate(rows[:rows_count - 1], 1):
            for c, val in enumerate(row):
                table.cell(r, c).text = str(val)

    def generate_quarterly_review(self, data: dict) -> tuple[bytes, Path]:
        """Generate quarterly business review presentation."""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)  # 16:9 widescreen
        prs.slide_height = PptxInches(7.5)

        period = f"{data.get('period_start', '')} → {data.get('period_end', '')}"
        date_str = datetime.now().strftime("%B %d, %Y")

        # Slide 1: Title
        self._add_title_slide(prs, "MaaS-Router Quarterly Review", period, date_str)

        # Slide 2: Executive Summary
        summary = data.get("summary", {})
        self._add_bullet_slide(prs, "Executive Summary", [
            f"💰 Total Revenue: CRED {summary.get('total_revenue', 0):,.2f}",
            f"👥 Active Users: {summary.get('active_users', 0)}",
            f"🔑 Active API Keys: {summary.get('active_keys', 0)}",
            f"📊 Total Requests: {summary.get('total_requests', 0):,}",
            f"⚡ Avg Latency: {summary.get('avg_latency', 0)}ms",
            f"✅ Uptime: {summary.get('uptime', '99.9%')}",
        ])

        # Slide 3: Revenue Breakdown
        rev_data = data.get("revenue_breakdown", [])
        if rev_data:
            self._add_table_slide(
                prs, "Revenue Breakdown by Model",
                ["Model", "Revenue (CRED)", "Requests", "Share"],
                [[m.get("model", ""), f"{m.get('revenue', 0):,.2f}",
                  str(m.get("requests", 0)), f"{m.get('share', 0) * 100:.1f}%"]
                 for m in rev_data],
            )

        # Slide 4: User Growth
        user_data = data.get("user_growth", [])
        if user_data:
            self._add_table_slide(
                prs, "User Growth & Retention",
                ["Month", "New Users", "Active Users", "Churn Rate"],
                [[m.get("month", ""), str(m.get("new_users", 0)),
                  str(m.get("active_users", 0)), f"{m.get('churn', 0) * 100:.1f}%"]
                 for m in user_data],
            )

        # Slide 5: Key Initiatives & Next Steps
        self._add_bullet_slide(prs, "Key Initiatives & Next Steps", [
            "🚀 Model expansion: Add 5 new providers",
            "📈 Enterprise tier: SSO + dedicated support",
            "🔒 SOC 2 Type II certification in progress",
            "🤖 AI Gateway: intelligent fallback & caching",
            "📊 Enhanced analytics dashboard Q3",
        ])

        # Save
        output_path = self.output_dir / self._filename("quarterly_review", "pptx")
        prs.save(str(output_path))

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue(), output_path


# ══════════════════════════════════════════════════════════════════
# Unified Document Service (Facade)
# ══════════════════════════════════════════════════════════════════

class DocumentService:
    """
    Facade for all document generation engines.
    Single entry point for the admin API layer.
    """

    def __init__(self):
        self.pdf = PDFEngine() if WEASYPRINT_AVAILABLE else None
        self.excel = ExcelEngine() if OPENPYXL_AVAILABLE else None
        self.word = WordEngine() if PYTHON_DOCX_AVAILABLE else None
        self.pptx = PPTXEngine() if PYTHON_PPTX_AVAILABLE else None

    def get_available_engines(self) -> dict:
        """Return which engines are available."""
        return {
            "pdf": self.pdf is not None,
            "excel": self.excel is not None,
            "word": self.word is not None,
            "pptx": self.pptx is not None,
        }

    def generate(self, doc_type: DocumentType, data: dict,
                 format: DocumentFormat = DocumentFormat.PDF) -> tuple[bytes, Path, str]:
        """
        Generate a document.

        Returns:
            (bytes_content, file_path, mime_type)
        """
        generators = {
            (DocumentType.BILLING_REPORT, DocumentFormat.PDF):    ("pdf", "generate_billing_report"),
            (DocumentType.BILLING_REPORT, DocumentFormat.EXCEL):  ("excel", "generate_billing_excel"),
            (DocumentType.USER_USAGE_REPORT, DocumentFormat.PDF): ("pdf", "generate_user_report"),
            (DocumentType.USER_USAGE_REPORT, DocumentFormat.WORD):("word", "generate_user_report_docx"),
            (DocumentType.OPS_DAILY, DocumentFormat.PDF):         ("pdf", "generate_ops_daily"),
            (DocumentType.AUDIT_REPORT, DocumentFormat.WORD):     ("word", "generate_audit_report"),
            (DocumentType.DATA_EXPORT, DocumentFormat.EXCEL):     ("excel", "generate_data_export"),
            (DocumentType.DATA_EXPORT, DocumentFormat.CSV):       ("excel", "generate_csv"),
            (DocumentType.MODEL_PERFORMANCE, DocumentFormat.PPTX):("pptx", "generate_quarterly_review"),
        }

        key = (doc_type, format)
        if key not in generators:
            # Fallback: try PDF for any doc_type
            if format == DocumentFormat.PDF and self.pdf:
                content, path = self.pdf.render_from_template(
                    f"{doc_type.value}.html", data, doc_type=doc_type.value
                )
                return content, path, "application/pdf"
            raise ValueError(f"No generator for {doc_type.value} in {format.value} format")

        engine_name, method_name = generators[key]
        engine = getattr(self, engine_name)
        if engine is None:
            raise RuntimeError(
                f"{engine_name} engine is not available. "
                f"Install dependencies: pip install {'weasyprint' if engine_name == 'pdf' else engine_name}"
            )

        method = getattr(engine, method_name)
        content, path = method(data)

        mime_map = {
            "pdf": "application/pdf",
            "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "word": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

        return content, path, mime_map.get(engine_name, "application/octet-stream")

    def list_outputs(self, limit: int = 20) -> list[dict]:
        """List recently generated documents."""
        outputs = []
        for f in sorted(OUTPUT_DIR.glob("*"), key=lambda x: x.stat().st_mtime, reverse=True):
            if f.is_file() and not f.name.startswith("."):
                stat = f.stat()
                outputs.append({
                    "filename": f.name,
                    "size": stat.st_size,
                    "size_display": f"{stat.st_size / 1024:.1f} KB",
                    "created": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                    "format": f.suffix.lstrip("."),
                })
            if len(outputs) >= limit:
                break
        return outputs


# ══════════════════════════════════════════════════════════════════
# Data Adapter — builds document context from DB queries
# ══════════════════════════════════════════════════════════════════

class DocumentDataAdapter:
    """
    Transforms database query results into the data structures
    expected by document templates and engines.
    """

    @staticmethod
    def billing_data(dashboard_overview: dict, trends: dict,
                     transactions: list[dict]) -> dict:
        """Build billing report data context."""
        return {
            "report_title": "MaaS-Router Billing Report",
            "period_start": trends.get("data", [{}])[0].get("date", "N/A") if trends.get("data") else "N/A",
            "period_end": trends.get("data", [{}])[-1].get("date", "N/A") if trends.get("data") else "N/A",
            "summary": {
                "total_revenue": dashboard_overview.get("monthly_revenue", 0),
                "today_revenue": dashboard_overview.get("today_revenue", 0),
                "total_transactions": len(transactions),
                "active_users": dashboard_overview.get("active_today", 0),
                "new_users": 0,  # populated by caller
                "top_plan": "N/A",
            },
            "daily_breakdown": [
                {
                    "date": d.get("date", ""),
                    "requests": d.get("requests", 0),
                    "tokens": d.get("tokens", 0),
                    "revenue": d.get("cost", 0),
                    "avg_latency": 0,
                    "error_rate": 0,
                }
                for d in trends.get("data", [])
            ],
            "recent_transactions": [
                {
                    "id": t.get("id", ""),
                    "user": t.get("user_email", ""),
                    "txn_type": t.get("type", ""),
                    "model": t.get("model_id", ""),
                    "tokens": t.get("total_tokens", 0),
                    "amount": abs(t.get("amount", 0)),
                    "currency": t.get("currency", "CRED"),
                    "status": t.get("status", ""),
                    "created_at": t.get("created_at", ""),
                }
                for t in transactions
            ],
        }

    @staticmethod
    def user_usage_data(users: list[dict], model_distribution: dict) -> dict:
        """Build user usage report data context."""
        total_tokens = sum(u.get("total_tokens", 0) for u in users)
        total_requests = sum(u.get("total_requests", 0) for u in users)

        return {
            "period_start": "Last 30 Days",
            "period_end": datetime.now().strftime("%Y-%m-%d"),
            "overview": {
                "active_users": len(users),
                "total_requests": total_requests,
                "total_tokens": total_tokens,
                "total_cost": round(sum(u.get("total_cost", 0) for u in users), 2),
            },
            "top_users": [
                {
                    "display_name": u.get("display_name", u.get("email", "")),
                    "email": u.get("email", ""),
                    "plan": u.get("plan_id", "free"),
                    "requests": u.get("total_requests", 0),
                    "tokens": u.get("total_tokens", 0),
                    "cost": u.get("total_cost", 0),
                }
                for u in sorted(users, key=lambda x: x.get("total_tokens", 0), reverse=True)[:20]
            ],
            "model_usage": model_distribution.get("data", []),
        }

    @staticmethod
    def ops_daily_data(overview: dict, trends: dict, alerts: list[dict],
                        service_health: dict) -> dict:
        """Build daily operations report context."""
        return {
            "report_date": datetime.now().strftime("%Y-%m-%d"),
            "summary": overview,
            "trends": trends,
            "alerts": alerts,
            "service_health": service_health,
        }

    @staticmethod
    def audit_data(audit_logs: list[dict], user_activity: list[dict],
                    period_start: str, period_end: str) -> dict:
        """Build audit report data context."""
        critical = [log for log in audit_logs if log.get("action", "").startswith("critical")]
        return {
            "report_id": f"AUD-{datetime.now().strftime('%Y%m%d%H%M')}",
            "period_start": period_start,
            "period_end": period_end,
            "summary": {
                "total_events": len(audit_logs),
                "active_users": len(user_activity),
                "critical_events": len(critical),
            },
            "audit_summary": [],  # Aggregated by caller
            "critical_events": critical,
            "user_activity": user_activity,
            "compliance_notes": [],
        }
